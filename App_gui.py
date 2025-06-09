import os, re, sys, threading, subprocess

import tkinter as tk
from tkinter import ttk, messagebox, filedialog, simpledialog

#import whisper
import fitz, pdfplumber
from pdfminer.high_level import extract_text
from pdfminer.layout import LAParams
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document


class FileProcessorApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("NBox v1.0")
        self.center_on_screen(350, 200)
        self.style = ttk.Style(self)
        self.style.theme_use('clam')
        self.style.configure("TButton", font=("Arial", 12), padding=10)
        self.progress = ttk.Progressbar(self, orient="horizontal", mode="determinate")
        self.progress.pack(side="bottom", fill="x")
        self.resizable(False, False)
        self.create_tabs()

    def center_on_screen(self, width, height):
        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()
        x = (screen_width // 2) - (width // 2)
        y = (screen_height // 2) - (height // 2)
        self.geometry(f"{width}x{height}+{x}+{y}")


    def center_window(self):
        self.update_idletasks()
        width = self.winfo_width()
        height = self.winfo_height()
        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()
        x = (screen_width // 2) - (width // 2)
        y = (screen_height // 2) - (height // 2)
        self.geometry(f"{width}x{height}+{x}+{y}")


    def create_tabs(self):
        
        notebook = ttk.Notebook(self)
        notebook.pack(expand=True, fill='both')
        
        tabs = [
            #("Речь из видео", self.transcribe_video_main), TO DO 
            #("Upscale img", self.upscale_img), TO DO 
            ("Сжатие видео", self.compress_video_main),
            ("PDF → TXT", self.pdf_to_txt_main),
            ("PDF → Word", self.pdf_to_word_main),
            ("PPTX → Word", self.pptx_to_word_main)
        ]
        
        for title, func in tabs:
            frame = ttk.Frame(notebook)
            notebook.add(frame, text=title)
            btn_frame = ttk.Frame(frame)
            btn_frame.pack(expand=True, fill='both', padx=10, pady=10)
            btn = ttk.Button(btn_frame, text=f"Выполнить: {title}", command=lambda f=func: threading.Thread(target=f).start())
            btn.pack(expand=True, fill='x')
        

        bottom_frame = ttk.Frame(self)
        bottom_frame.pack(side='bottom', fill='x', padx=10, pady=5)

        about_button = ttk.Button(bottom_frame, text="ℹ️ О программе", command=self.show_about)
        about_button.pack(side='left', expand=True, fill='x', padx=(0, 5))

        instructions_button = ttk.Button(bottom_frame, text="📖 Инструкция", command=self.show_instructions)
        instructions_button.pack(side='right', expand=True, fill='x', padx=(5, 0))

    def show_instructions(self):
        window = tk.Toplevel(self)
        window.title("📖 Инструкция по использованию")
        window.resizable(False, False)

        width, height = 700, 500
        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()
        x = (screen_width // 2) - (width // 2)
        y = (screen_height // 2) - (height // 2)
        window.geometry(f"{width}x{height}+{x}+{y}")

        text_widget = tk.Text(window, wrap='word', font=("Arial", 11))
        scrollbar = ttk.Scrollbar(window, orient="vertical", command=text_widget.yview)
        text_widget.configure(yscrollcommand=scrollbar.set)

        text_widget.pack(side="left", fill="both", expand=True, padx=(10, 0), pady=10)
        scrollbar.pack(side="right", fill="y", pady=10, padx=(0, 10))

        instructions = (
            "Инструкция по использованию программы NBox\n\n"
            "Шаг 1. Запуск программы\n"
            "1. Откройте папку с программой NBox.\n"
            "2. Дважды кликните по файлу NBox.exe\n"
            "3. Откроется главное окно с вкладками:\n"
            "   Сжатие видео, PDF → TXT, PDF → Word, PPTX → Word.\n\n"
            "Шаг 2. Сжатие видео\n"
            "1. Перейдите во вкладку «Сжатие видео».\n"
            "2. Нажмите кнопку «Выполнить: Сжатие видео».\n"
            "3. Выберите один или несколько видеофайлов (mp4, mov, avi, mkv, webm).\n"
            "4. В появившемся окне выберите качество сжатия:\n"
            "   - Максимальное качество — 3000k\n"
            "   - Среднее качество — 1500k\n"
            "   - Сильное сжатие — 800k\n"
            "5. После сжатия файлы будут сохранены рядом с оригиналом с суффиксом _compressed.\n\n"
            "Шаг 3. Конвертация PDF в TXT\n"
            "1. Перейдите во вкладку «PDF → TXT».\n"
            "2. Нажмите кнопку «Выполнить: PDF → TXT».\n"
            "3. Выберите один или несколько PDF-файлов.\n"
            "4. Программа извлечёт текст и изображения:\n"
            "   - Текст будет сохранён в .txt\n"
            "   - Изображения — в отдельную папку рядом с исходником\n"
            "5. По завершении появится сообщение «PDF успешно конвертированы в TXT с изображениями!»\n\n"
            "Шаг 4. Конвертация PDF в Word\n"
            "1. Перейдите во вкладку «PDF → Word».\n"
            "2. Нажмите кнопку «Выполнить: PDF → Word».\n"
            "3. Выберите PDF-файлы.\n"
            "4. Будет создан .docx файл с извлечённым текстом и ссылками на изображения.\n\n"
            "Шаг 5. Конвертация PPTX в Word\n"
            "1. Перейдите во вкладку «PPTX → Word».\n"
            "2. Нажмите кнопку «Выполнить: PPTX → Word».\n"
            "3. Выберите презентации (.pptx).\n"
            "4. Программа извлечёт:\n"
            "   - Текст слайдов\n"
            "   - Заметки к слайдам\n"
            "   - Изображения (сохранит отдельно и укажет в документе)\n\n"
            "• Все функции работают локально, без интернета.\n"
            "• Вы можете запустить несколько операций параллельно (например, сжатие и конвертацию).\n"
            "• Используйте кнопку «Закрыть» в правом верхнем углу для выхода.\n"
        )

        text_widget.insert("1.0", instructions)
        text_widget.config(state="disabled")  # запрет на редактирование


    def show_about(self):
        messagebox.showinfo("О программе", "Автор: Exodia\nЛицензия: Только для личного использования\nНе использовать в контуре!")

    def update_progress(self, value, maximum):
        self.progress["value"] = value
        self.progress["maximum"] = maximum
        self.update_idletasks()
    #TO DO - Пересобрать уже с транскрибацией 
    # def transcribe_video_main(self):
    #     video_path = filedialog.askopenfilename(filetypes=[("Видео", "*.mp4 *.mov *.avi *.mkv *.webm")])
    #     if not video_path:
    #         return
    #     audio_path = "temp_audio.wav"
    #     self.update_progress(0, 100)
    #     subprocess.run(["ffmpeg", "-y", "-i", video_path, "-ar", "16000", "-ac", "1", "-f", "wav", audio_path], check=True)
    #     self.update_progress(20, 100)
    #     model = whisper.load_model("small")
    #     result = model.transcribe(audio_path, language="ru")
    #     segments = result.get("segments", [])
    #     total = len(segments)
    #     text = ""
    #     for i, seg in enumerate(segments):
    #         text += seg["text"] + " "
    #         self.update_progress(20 + (i + 1) * 80 // total, 100)
    #     txt_path = os.path.splitext(video_path)[0] + ".txt"
    #     with open(txt_path, "w", encoding="utf-8") as f:
    #         f.write(text)
    #     os.remove(audio_path)
    #     messagebox.showinfo("Готово", f"Текст сохранён в: {txt_path}")
    #     self.update_progress(0, 100)  
    
    def compress_video_main(self):
        files = filedialog.askopenfilenames(filetypes=[("Видео", "*.mp4 *.mov *.avi *.mkv *.webm")])
        if not files:
            return
        self.get_bitrate_choice_and_compress(files)
    
    def get_bitrate_choice_and_compress(self, files):
        choices = {"Максимальное качество": "3000k", "Среднее качество": "1500k", "Сильное сжатие": "800k"}
        choice_window = tk.Toplevel(self)
        choice_window.title("Выбор битрейта")
        choice_window.geometry(f"300x200+{self.winfo_x() + 50}+{self.winfo_y() + 50}")
        choice_window.transient(self)
        choice_window.grab_set()
        tk.Label(choice_window, text="Выберите степень сжатия:").pack(padx=10, pady=10)
        for label, value in choices.items():
            btn = ttk.Button(choice_window, text=label, command=lambda v=value: self.set_bitrate_and_compress(choice_window, files, v))
            btn.pack(padx=10, pady=5, fill='x')

    def get_ffmpeg_path(self):
        if getattr(sys, 'frozen', False):
            base_path = sys._MEIPASS
            return os.path.join(base_path, "bin", "ffmpeg.exe")
        else:
            base_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bin")
            return os.path.join(base_path, "ffmpeg.exe")


    def set_bitrate_and_compress(self, window, files, bitrate):
        window.destroy()
        total = len(files)
        ffmpeg_path = self.get_ffmpeg_path()  
        for i, filepath in enumerate(files):
            output_path = os.path.splitext(filepath)[0] + f"_compressed.mp4"
            subprocess.run([ffmpeg_path, "-i", filepath, "-b:v", bitrate, "-c:a", "aac", "-y", output_path], check=True)
            self.update_progress(i + 1, total)
        messagebox.showinfo("Готово", "Все видео успешно сжаты!")
        self.update_progress(0, total)

    def pdf_to_txt_main(self):
        files = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if not files:
            return
        total = len(files)
        for i, filepath in enumerate(files):
            full_text = ''
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        full_text += text + '\n\n'
            cleaned_text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', full_text.strip())
            txt_path = os.path.splitext(filepath)[0] + ".txt"
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(cleaned_text)
            
            img_folder = os.path.splitext(filepath)[0] + "_images"
            os.makedirs(img_folder, exist_ok=True)
            pdf_file = fitz.open(filepath)
            for page_number in range(len(pdf_file)):
                page = pdf_file[page_number]
                images = page.get_images(full=True)
                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = pdf_file.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    img_filename = f"page{page_number+1}_img{img_index+1}.{image_ext}"
                    img_path = os.path.join(img_folder, img_filename)
                    with open(img_path, 'wb') as img_file:
                        img_file.write(image_bytes)
            pdf_file.close()
            
            self.update_progress(i + 1, total)
        messagebox.showinfo("Готово", "PDF успешно конвертированы в TXT с изображениями!")
        self.update_progress(0, total)

    def pdf_to_word_main(self):
        files = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if not files:
            return
        total = len(files)
        for i, filepath in enumerate(files):
            full_text = ''
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        full_text += text + '\n\n'
            output_path = os.path.splitext(filepath)[0] + ".docx"
            doc = Document()
            paragraphs = re.split(r'\n{2,}', full_text)
            for paragraph in paragraphs:
                paragraph = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', paragraph.strip())
                if paragraph:
                    doc.add_paragraph(paragraph)
            
            img_folder = os.path.splitext(filepath)[0] + "_images"
            os.makedirs(img_folder, exist_ok=True)
            pdf_file = fitz.open(filepath)
            for page_number in range(len(pdf_file)):
                page = pdf_file[page_number]
                images = page.get_images(full=True)
                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = pdf_file.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    img_filename = f"page{page_number+1}_img{img_index+1}.{image_ext}"
                    img_path = os.path.join(img_folder, img_filename)
                    with open(img_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    doc.add_paragraph(f"[Изображение сохранено: {img_filename}]")
            pdf_file.close()
            
            doc.save(output_path)
            self.update_progress(i + 1, total)
        messagebox.showinfo("Готово", "PDF успешно конвертированы в Word с изображениями!")
        self.update_progress(0, total)
    
    def pptx_to_word_main(self):
        files = filedialog.askopenfilenames(filetypes=[("PPTX", "*.pptx")])
        if not files:
            return
        total = len(files)
        for filepath in files:
            prs = Presentation(filepath)
            doc = Document()
            doc.add_heading("Презентация: " + os.path.basename(filepath), 0)
            img_folder = os.path.splitext(filepath)[0] + "_images"
            os.makedirs(img_folder, exist_ok=True)
            for i, slide in enumerate(prs.slides):
                doc.add_heading(f"Слайд {i+1}", level=1)
                pic_count = 1
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = ' '.join(run.text.strip() for para in shape.text_frame.paragraphs for run in para.runs if run.text.strip())
                        if text:
                            doc.add_paragraph(text)
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        ext = image.ext
                        image_bytes = image.blob
                        img_filename = f"slide{i+1}_pic{pic_count}.{ext}"
                        img_path = os.path.join(img_folder, img_filename)
                        with open(img_path, 'wb') as img_file:
                            img_file.write(image_bytes)
                        pic_count += 1
                        doc.add_paragraph(f"[Изображение сохранено: {img_filename}]")
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        cleaned_notes = '\n'.join(line.strip() for line in notes.splitlines() if line.strip())
                        doc.add_paragraph("Заметки к слайду:")
                        doc.add_paragraph(cleaned_notes)
            output_path = os.path.splitext(filepath)[0] + ".docx"
            doc.save(output_path)
            self.update_progress(files.index(filepath)+1, total)
        messagebox.showinfo("Готово", "Презентации успешно сконвертированы в Word с изображениями!")
        self.update_progress(0, total)

if __name__ == "__main__":
    app = FileProcessorApp()
    app.mainloop()
